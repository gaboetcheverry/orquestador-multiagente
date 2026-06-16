import streamlit as st
import os
import html  # ← IMPORTACIÓN AGREGADA
from datetime import datetime
from docx import Document
from pptx import Presentation
import plotly.express as px
import pandas as pd

# Importaciones para PDF Profesional
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
from reportlab.lib.colors import HexColor

# Importaciones para LLM Robusto
from langchain_openai import ChatOpenAI
from langchain_core.prompts import PromptTemplate

st.set_page_config(page_title="Orquestador Estratégico Multiagente Pro", layout="wide")

# ==============================================================================
# 1. CONFIGURACIÓN DE AGENTES Y PROMPTS ROBUSTOS
# ==============================================================================
AGENTES_DOMINIO = [
    "Agente Financiero (ROI, CAPEX, OPEX, Flujo de Caja)",
    "Agente Macroeconómico (Inflación, Tasas, Geopolítica)",
    "Agente de Riesgos (Mitigación, Escenarios, Probabilidades)",
    "Agente de Sistemas (Infraestructura, Ciberseguridad, Legacy)",
    "Agente de Talento (Cultura, Retención, Upskilling)",
    "Agente ESG (Sostenibilidad, Gobernanza, Impacto Social)",
    "Agente de Mercado (Competencia, Cuota, Tendencias de Consumo)",
    "Agente Digital (Transformación, Automatización, Data)",
    "Agente de Innovación (I+D, Nuevos Modelos de Negocio)",
    "Agente de Prospectiva (Escenarios a 5 y 10 años)",
    "Agente Probabilístico (Análisis Bayesiano, Montecarlo)"
]

def ejecutar_cadena_multiagente(tema_crudo: str, api_key: str):
    """Ejecuta la orquestación real usando LLM para cada etapa."""
    os.environ["OPENAI_API_KEY"] = api_key
    # Usamos gpt-4o-mini por ser más rápido y económico, pero igual de robusto para este tipo de tareas. 
    # Si prefieres el completo, cambia a "gpt-4o"
    llm = ChatOpenAI(model="gpt-4o-mini", temperature=0.3)
    
    progress_bar = st.progress(0.0)
    status_text = st.empty()
    
    resultados = {"tema_original": tema_crudo, "fecha": datetime.now().strftime("%d de %B de %Y")}
    
    # Total de pasos para la barra de progreso (1 briefing + 11 agentes + 1 auditor + 1 consenso = 14)
    total_steps = 14
    current_step = 0

    # PASO 1: Agente de Limpieza y Refinamiento
    current_step += 1
    status_text.text("🧹 Agente de Limpieza: Estructurando el briefing estratégico...")
    progress_bar.progress(current_step / total_steps)
    
    prompt_refinar = PromptTemplate.from_template(
        "Eres un Consultor Estratégico Senior. El usuario dio esta consulta cruda: '{tema}'. "
        "Reescríbela como un 'Briefing Estratégico Ejecutivo' formal, detallando: Objetivo Central, Alcance, Restricciones Asumidas y KPIs de Éxito. Sé conciso pero muy profesional."
    )
    briefing = llm.invoke(prompt_refinar.format(tema=tema_crudo)).content
    resultados["briefing"] = briefing

    # PASO 2: Análisis de los 11 Agentes Especializados
    resultados["analisis_dominios"] = {}
    for i, agente in enumerate(AGENTES_DOMINIO):
        current_step += 1
        status_text.text(f"🧠 {agente} generando análisis profundo... ({i+1}/11)")
        progress_bar.progress(current_step / total_steps)
        
        prompt_agente = PromptTemplate.from_template(
            "Eres el {rol}. Analiza el siguiente Briefing Estratégico: \n\n'{briefing}'\n\n"
            "Genera un análisis robusto, técnico y accionable (mínimo 3 párrafos). Incluye: "
            "1. Diagnóstico actual, 2. Riesgos/Oportunidades específicos de tu dominio, 3. Recomendación táctica inmediata."
        )
        analisis = llm.invoke(prompt_agente.format(rol=agente, briefing=briefing)).content
        resultados["analisis_dominios"][agente] = analisis

    # PASO 3: Auditor Crítico
    current_step += 1
    status_text.text("⚖️ Auditor Crítico: Buscando sesgos, contradicciones y viabilidad...")
    progress_bar.progress(current_step / total_steps)
    
    prompt_auditor = PromptTemplate.from_template(
        "Eres un Auditor Crítico implacable. Revisa los análisis de los 11 agentes sobre este briefing: '{briefing}'. "
        "Identifica: 1. Contradicciones entre agentes (ej: Finanzas vs ESG), 2. Suposiciones no validadas, 3. Puntos ciegos críticos. "
        "Emite un veredicto de viabilidad (Alta, Media, Baja) con justificación."
    )
    auditoria = llm.invoke(prompt_auditor.format(briefing=briefing)).content
    resultados["auditoria"] = auditoria

    # PASO 4: Motor de Consenso
    current_step += 1
    status_text.text("🤝 Motor de Consenso: Sintetizando la estrategia final...")
    progress_bar.progress(current_step / total_steps) # Esto ahora es exactamente 1.0 (14/14)
    
    prompt_consenso = PromptTemplate.from_template(
        "Eres el CEO y Motor de Consenso. Basado en el Briefing: '{briefing}', los 11 análisis y la Auditoría Crítica: '{auditoria}', "
        "redacta la 'Estrategia Unificada Final'. Debe incluir: Resumen Ejecutivo (1 párrafo), 3 Pilares Estratégicos Prioritarios, "
        "Hoja de Ruta a 90 días y Métrica de Éxito principal. Sé directo y orientado a la acción."
    )
    consenso = llm.invoke(prompt_consenso.format(briefing=briefing, auditoria=auditoria)).content
    resultados["consenso"] = consenso
    
    status_text.text("✅ ¡Orquestación Multiagente Completada con Éxito!")
    progress_bar.progress(1.0)
    
    return resultados

# ==============================================================================
# 2. GENERADOR DE PDF PROFESIONAL (PLATYPUS)
# ==============================================================================
def generar_pdf_profesional(data, filename="Reporte_Estrategico_Pro.pdf"):
    doc = SimpleDocTemplate(filename, pagesize=A4, rightMargin=50, leftMargin=50, topMargin=50, bottomMargin=50)
    story = []
    styles = getSampleStyleSheet()
    
    # Estilos personalizados
    styles.add(ParagraphStyle(name='TituloPrincipal', fontName='Helvetica-Bold', fontSize=22, textColor=HexColor('#1f3a93'), alignment=TA_CENTER, spaceAfter=30))
    styles.add(ParagraphStyle(name='Subtitulo', fontName='Helvetica-Bold', fontSize=14, textColor=HexColor('#1f3a93'), spaceBefore=20, spaceAfter=10))
    styles.add(ParagraphStyle(name='Cuerpo', fontName='Helvetica', fontSize=11, leading=16, alignment=TA_JUSTIFY, spaceAfter=12))
    styles.add(ParagraphStyle(name='CajaResumen', fontName='Helvetica', fontSize=11, leading=16, backColor=HexColor('#f0f4f8'), borderColor=HexColor('#1f3a93'), borderWidth=1, borderPadding=10, spaceAfter=20))

    # Portada / Encabezado
    story.append(Paragraph("REPORTE ESTRATÉGICO MULTIAGENTE", styles['TituloPrincipal']))
    story.append(Paragraph(f"Tema: {html.escape(data['tema_original'])}", styles['Subtitulo']))
    story.append(Paragraph(f"Fecha de Emisión: {data['fecha']}", styles['Cuerpo']))
    story.append(Spacer(1, 30))

    # Consenso (Lo más importante primero)
    story.append(Paragraph("1. CONSENSO ESTRATÉGICO Y HOJA DE RUTA", styles['Subtitulo']))
    # ← CORRECCIÓN: Escapar primero, luego reemplazar saltos de línea
    texto_consenso_seguro = html.escape(str(data['consenso']))
    texto_consenso_formateado = texto_consenso_seguro.replace('\r\n', '<br/>').replace('\n', '<br/>')
    story.append(Paragraph(texto_consenso_formateado, styles['CajaResumen']))
    story.append(PageBreak())

    # Auditoría
    story.append(Paragraph("2. DICTAMEN DEL AUDITOR CRÍTICO", styles['Subtitulo']))
    # ← CORRECCIÓN: Escapar primero, luego reemplazar saltos de línea
    texto_auditoria_seguro = html.escape(str(data['auditoria']))
    texto_auditoria_formateado = texto_auditoria_seguro.replace('\r\n', '<br/>').replace('\n', '<br/>')
    story.append(Paragraph(texto_auditoria_formateado, styles['Cuerpo']))
    story.append(Spacer(1, 20))

    # Análisis por Dominio
    story.append(PageBreak())
    story.append(Paragraph("3. ANÁLISIS DETALLADO POR DOMINIO DE EXPERTOS", styles['Subtitulo']))
    
    for agente, analisis in data['analisis_dominios'].items():
        story.append(Paragraph(html.escape(agente), styles['Subtitulo']))
        # ← CORRECCIÓN: Escapar primero, luego reemplazar saltos de línea
        texto_analisis_seguro = html.escape(str(analisis))
        texto_analisis_formateado = texto_analisis_seguro.replace('\r\n', '<br/>').replace('\n', '<br/>')
        story.append(Paragraph(texto_analisis_formateado, styles['Cuerpo']))
        story.append(Spacer(1, 15))

    doc.build(story)
    return filename

# ==============================================================================
# 3. GENERADORES WORD Y PPTX
# ==============================================================================
def generar_word(data, filename="Reporte_Estrategico.docx"):
    doc = Document()
    doc.add_heading(f'Reporte Estratégico: {data["tema_original"]}', 0)
    doc.add_paragraph(f'Fecha: {data["fecha"]}')
    doc.add_heading('Consenso Ejecutivo', level=1)
    doc.add_paragraph(data['consenso'])
    doc.add_heading('Auditoría Crítica', level=1)
    doc.add_paragraph(data['auditoria'])
    for agente, analisis in data['analisis_dominios'].items():
        doc.add_heading(agente, level=2)
        doc.add_paragraph(analisis)
    doc.save(filename)
    return filename

def generar_pptx(data, filename="Presentacion_Estrategica.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Estrategia: {data['tema_original']}"
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Consenso del Motor Estratégico"
    slide.placeholders[1].text = data['consenso'][:1000] + "..."
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Hallazgos del Auditor Crítico"
    slide.placeholders[1].text = data['auditoria'][:1000] + "..."
    prs.save(filename)
    return filename

# ==============================================================================
# 4. INTERFAZ DE USUARIO (STREAMLIT)
# ==============================================================================
def main():
    st.title("🏢 Orquestador Estratégico Multiagente PRO")
    st.markdown("Sistema de inteligencia colectiva con **Agente de Limpieza**, 11 Especialistas, Auditoría Crítica y Motor de Consenso.")
    
    with st.sidebar:
        st.header("⚙️ Configuración del Motor")
        api_key = st.text_input("🔑 OpenAI API Key", type="password", value=os.getenv("OPENAI_API_KEY", ""))
        st.markdown("---")
        st.info("Arquitectura Activa:\n1. Refinador de Consultas\n2. 11 Agentes de Dominio\n3. Auditor Crítico\n4. Motor de Consenso")

    st.markdown("### 🎯 Definir Objetivo Estratégico")
    tema = st.text_area("Describe el proyecto, empresa o desafío a analizar:", 
                        placeholder="Ej: Evaluar la viabilidad de expandir nuestra planta de manufactura a Vietnam en 2025, considerando riesgos geopolíticos, costos laborales y cumplimiento ESG.", height=100)
    
    if st.button("🚀 INICIAR ORQUESTACIÓN MULTIAGENTE", type="primary", use_container_width=True):
        if not tema:
            st.warning("Por favor, ingrese un tema para analizar.")
            return
        if not api_key:
            st.error("⚠️ Es obligatorio proporcionar una OpenAI API Key válida para el análisis robusto.")
            return
            
        with st.spinner("Procesando cadena multiagente... (Esto toma 1-2 minutos)"):
            try:
                resultados = ejecutar_cadena_multiagente(tema, api_key)
                st.session_state['resultados'] = resultados
                st.success("✅ Análisis multiagente completado exitosamente.")
            except Exception as e:
                st.error(f"Error en la API: {str(e)}. Verifica tu clave.")
                return

    if 'resultados' in st.session_state:
        data = st.session_state['resultados']
        
        st.markdown("---")
        st.subheader("📊 Dashboard de Consenso Estratégico")
        
        col1, col2 = st.columns([2, 1])
        with col1:
            st.markdown("#### 🎯 Estrategia Unificada (Motor de Consenso)")
            st.info(data['consenso'])
        with col2:
            st.markdown("#### ⚖️ Veredicto del Auditor")
            st.warning(data['auditoria'][:400] + "...")
            
        st.markdown("#### 🧠 Briefing Estratégico (Generado por Agente de Limpieza)")
        st.markdown(f"*{data['briefing']}*")

        st.markdown("---")
        st.subheader("📥 Centro de Descarga de Reportes Profesionales")
        
        colA, colB, colC = st.columns(3)
        
        pdf_file = generar_pdf_profesional(data)
        word_file = generar_word(data)
        pptx_file = generar_pptx(data)
        
        with colA:
            with open(pdf_file, "rb") as f:
                st.download_button(label="📄 Descargar PDF Profesional", data=f, file_name=pdf_file, mime="application/pdf", use_container_width=True)
        with colB:
            with open(word_file, "rb") as f:
                st.download_button(label="📝 Descargar Word", data=f, file_name=word_file, mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
        with colC:
            with open(pptx_file, "rb") as f:
                st.download_button(label="📊 Descargar PowerPoint", data=f, file_name=pptx_file, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)

if __name__ == "__main__":
    main()
